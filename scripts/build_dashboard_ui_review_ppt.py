"""Build the dashboard UI/UX redesign management review deck.

The deck uses only the published audit conclusions and high-fidelity prototype
screenshots. It does not read production data or call application interfaces.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
AUDIT_PATH = REPORTS / "2026-07-20-dashboard-ui-ux-audit.md"
OUTPUT_PATH = REPORTS / "2026-07-20-dashboard-ui-ux-redesign.pptx"
SCREENSHOT_DIR = REPORTS / "dashboard-ui-demo" / "screenshots"
SCREENSHOTS = {
    "overview": SCREENSHOT_DIR / "overview-desktop.png",
    "automation": SCREENSHOT_DIR / "automation-desktop.png",
    "analytics": SCREENSHOT_DIR / "analytics-desktop.png",
    "records": SCREENSHOT_DIR / "records-desktop.png",
    "settings": SCREENSHOT_DIR / "settings-desktop.png",
    "records_mobile": SCREENSHOT_DIR / "records-mobile.png",
}

BG = "07131D"
PANEL = "0D2230"
PANEL_2 = "123143"
TEXT = "F5F8FA"
MUTED = "8EA8B5"
CYAN = "36D7C4"
BLUE = "63AFFF"
GREEN = "68D391"
ORANGE = "FFB454"
RED = "FF6B6B"
PURPLE = "A995FF"
YELLOW = "F5D76E"
WHITE = "FFFFFF"
FONT = "Microsoft YaHei"
FONT_NUM = "Aptos Display"
FOOTER_NOTE = "高保真演示原型，不连接真实接口"


def _load_pptx() -> None:
    """Load optional deck dependencies only after input validation succeeds."""
    global RGBColor, MSO_CONNECTOR, MSO_SHAPE, MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    global Presentation, Inches, Pt
    from pptx import Presentation as _Presentation
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR as _MSO_CONNECTOR, MSO_SHAPE as _MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR, MSO_AUTO_SIZE as _MSO_AUTO_SIZE, PP_ALIGN as _PP_ALIGN
    from pptx.util import Inches as _Inches, Pt as _Pt

    Presentation = _Presentation
    RGBColor = _RGBColor
    MSO_CONNECTOR = _MSO_CONNECTOR
    MSO_SHAPE = _MSO_SHAPE
    MSO_ANCHOR = _MSO_ANCHOR
    MSO_AUTO_SIZE = _MSO_AUTO_SIZE
    PP_ALIGN = _PP_ALIGN
    Inches = _Inches
    Pt = _Pt


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_box(slide, left, top, width, height, fill=PANEL, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.8)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, TypeError):
            pass
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=16,
    color=TEXT,
    bold=False,
    font=FONT,
    align=None,
    valign=None,
):
    align = PP_ALIGN.LEFT if align is None else align
    valign = MSO_ANCHOR.TOP if valign is None else valign
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Pt(3)
    frame.margin_right = Pt(3)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines, left, top, width, height, size=13, gap=5):
    """Add short colored lines without relying on PowerPoint bullet rendering."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Pt(2)
    frame.margin_right = Pt(2)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    for index, (text, color, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_title(slide, eyebrow, title, page):
    add_text(slide, eyebrow.upper(), 0.70, 0.33, 5.7, 0.25, 9, CYAN, True)
    add_text(slide, title, 0.67, 0.66, 11.8, 0.60, 25, TEXT, True)
    add_text(
        slide,
        f"{page:02d} / 11",
        11.86,
        0.38,
        0.72,
        0.22,
        9,
        MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide, source="来源：UI/UX 审计报告与演示原型"):
    add_text(slide, source, 0.70, 7.15, 8.5, 0.16, 7.5, MUTED)
    add_text(
        slide,
        FOOTER_NOTE,
        9.15,
        7.15,
        3.45,
        0.16,
        7.5,
        MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_background(slide):
    add_box(slide, 0, 0, 13.333, 7.5, BG, BG)
    add_box(slide, 0.0, 0.0, 0.07, 7.5, CYAN, CYAN)
    for x, color in ((9.85, PANEL_2), (10.45, PANEL_2), (11.05, PANEL_2)):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x),
            Inches(0),
            Inches(x + 1.8),
            Inches(1.8),
        )
        line.line.color.rgb = rgb(color)
        line.line.width = Pt(0.7)


def add_badge(slide, text, left, top, width, color):
    add_box(slide, left, top, width, 0.30, PANEL_2, color, True)
    add_text(
        slide,
        text,
        left + 0.05,
        top + 0.025,
        width - 0.10,
        0.22,
        9,
        color,
        True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_picture_crop(slide, path: Path, left, top, width, height):
    """Add a screenshot to a stable frame without extra image dependencies."""
    return slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width), Inches(height)
    )


def add_screenshot_frame(slide, path: Path, left, top, width, height, label):
    add_box(slide, left - 0.03, top - 0.03, width + 0.06, height + 0.06, PANEL, BLUE, True)
    add_picture_crop(slide, path, left, top, width, height)
    add_badge(slide, label, left + 0.18, top + 0.16, 1.70, CYAN)


def add_callout(slide, number, title, body, left, top, width, color):
    add_text(slide, f"{number:02d}", left, top, 0.45, 0.32, 14, color, True, FONT_NUM)
    add_text(slide, title, left + 0.50, top - 0.01, width - 0.50, 0.28, 13, TEXT, True)
    add_text(slide, body, left + 0.50, top + 0.34, width - 0.50, 0.60, 10.5, MUTED)


def validate_inputs() -> str:
    if not AUDIT_PATH.is_file():
        raise FileNotFoundError(f"缺少审计报告：{AUDIT_PATH}")
    audit_text = AUDIT_PATH.read_text(encoding="utf-8")
    required_headings = (
        "P0：错误被伪装成“空数据”",
        "P0：认证取消会抑制后续认证",
        "P1：单页承担六类高频任务",
        "P2：字号与非文本对比度",
        "实施优先级",
    )
    missing_headings = [item for item in required_headings if item not in audit_text]
    if missing_headings:
        raise ValueError("审计报告缺少预期章节：" + "、".join(missing_headings))
    missing_screenshots = [str(path) for path in SCREENSHOTS.values() if not path.is_file()]
    if missing_screenshots:
        formatted = "\n  - ".join(missing_screenshots)
        raise FileNotFoundError(
            "无法生成 PPT：以下高保真截图缺失，请先生成截图后重试：\n  - " + formatted
        )
    return audit_text


def build_deck(output: Path) -> None:
    # Keep optional presentation dependencies out of the validation path so a
    # missing screenshot reports the actionable repository error first.
    _load_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 01 / Cover
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_text(slide, "GOODJOB  ·  管理评审", 0.78, 0.70, 4.5, 0.28, 11, CYAN, True)
    add_text(slide, "仪表盘 UI/UX\n重设计评审", 0.72, 1.43, 6.15, 1.55, 38, TEXT, True)
    add_text(
        slide,
        "从“单页堆叠”转向五个独立工作区，\n优先保障事实可信、控制可确认、任务可恢复。",
        0.78,
        3.35,
        5.80,
        0.92,
        17,
        MUTED,
    )
    add_box(slide, 7.35, 1.05, 5.10, 4.98, PANEL, PANEL_2, True)
    add_text(slide, "评审目标", 7.77, 1.45, 2.0, 0.30, 12, CYAN, True)
    cover_items = [
        ("01", "建立可信状态", "错误、未授权与真空数据明确分离", RED),
        ("02", "拆分工作区", "总览 / 自动化 / 分析 / 记录 / 设置", BLUE),
        ("03", "降低操作风险", "命令目标、执行状态、回执结果分层", GREEN),
    ]
    for index, (number, title, body, color) in enumerate(cover_items):
        y = 2.02 + index * 1.10
        add_text(slide, number, 7.78, y, 0.46, 0.34, 16, color, True, FONT_NUM)
        add_text(slide, title, 8.34, y - 0.02, 3.45, 0.28, 14, TEXT, True)
        add_text(slide, body, 8.34, y + 0.34, 3.55, 0.34, 10.5, MUTED)
    add_text(slide, "2026.07.20", 0.78, 6.52, 2.0, 0.30, 12, ORANGE, True, FONT_NUM)
    add_footer(slide)

    # 02 / Baseline
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "01 / 现状基线", "已有基础可复用，但信息结构与状态表达正在透支信任", 2)
    strengths = [
        ("回执链路", "连接、执行、同步已分轴，具备控制反馈基础", GREEN),
        ("渐进配置", "计划区已有展开、输入与焦点恢复机制", CYAN),
        ("响应式基础", "920 / 680 / 430px 已有分层重排", BLUE),
        ("可访问起点", "主导航、关键指标与确认操作已有语义", PURPLE),
    ]
    for index, (title, body, color) in enumerate(strengths):
        x = 0.75 + (index % 2) * 3.12
        y = 1.72 + (index // 2) * 1.48
        add_box(slide, x, y, 2.82, 1.18, PANEL, color, True)
        add_text(slide, title, x + 0.20, y + 0.18, 2.35, 0.26, 13, color, True)
        add_text(slide, body, x + 0.20, y + 0.55, 2.35, 0.42, 10.5, MUTED)
    add_box(slide, 7.08, 1.72, 5.48, 4.68, "131F2B", RED, True)
    add_text(slide, "当前管理风险", 7.48, 2.10, 2.5, 0.30, 14, RED, True)
    risks = [
        ("事实风险", "请求失败、未授权与零记录落入同一视觉结果"),
        ("操作风险", "命令已提交与浏览器实际生效缺少清晰边界"),
        ("效率风险", "六类任务共享长页、全局刷新与滚动定位"),
        ("审计风险", "筛选分子与全量分母混合，趋势可能不可比"),
    ]
    for index, (title, body) in enumerate(risks):
        y = 2.75 + index * 0.76
        add_text(slide, "!", 7.50, y, 0.30, 0.30, 15, RED, True, FONT_NUM, PP_ALIGN.CENTER)
        add_text(slide, title, 7.92, y, 1.05, 0.25, 11.5, TEXT, True)
        add_text(slide, body, 8.98, y, 3.10, 0.43, 10.2, MUTED)
    add_text(
        slide,
        "方向：保留可靠的控制与响应式基础，重构信息架构、请求状态机和高风险编辑路径。",
        0.82,
        6.58,
        11.70,
        0.34,
        14,
        ORANGE,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "来源：审计报告“现状优势”与“主要问题”")

    # 03 / Core issues
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "02 / 核心问题", "优先级不是视觉偏好，而是业务事实与操作风险", 3)
    issue_columns = [
        (
            "P0 事实可信",
            RED,
            [
                "错误被伪装成空数据",
                "认证取消后缺少恢复入口",
                "401 / 网络 / 5xx 未清晰区分",
            ],
            "先修：状态机 + 重试 + 重新验证",
        ),
        (
            "P1 运营风险",
            ORANGE,
            [
                "六类任务挤在同一长页",
                "期望 / 执行 / 回执语义混淆",
                "异步编辑可能被轮询覆盖",
                "弹窗与抽屉焦点链路不完整",
            ],
            "再拆：工作区 + 状态模型 + 并发策略",
        ),
        (
            "P2 效率与可达",
            BLUE,
            [
                "8-10px 字号与低对比文本",
                "ARIA / 键盘模式不一致",
                "移动端仍保留高密表格思路",
                "漏斗统计口径不可稳定比较",
            ],
            "后提：可读性 + 语义 + 统计契约",
        ),
    ]
    for index, (heading, color, bullets, action) in enumerate(issue_columns):
        x = 0.72 + index * 4.18
        add_box(slide, x, 1.66, 3.82, 4.70, PANEL, color, True)
        add_text(slide, heading, x + 0.26, 1.96, 3.22, 0.34, 16, color, True)
        add_rich_lines(
            slide,
            [(f"• {item}", TEXT, False) for item in bullets],
            x + 0.28,
            2.62,
            3.20,
            2.40,
            12,
            8,
        )
        add_box(slide, x + 0.24, 5.45, 3.34, 0.58, PANEL_2, color, True)
        add_text(slide, action, x + 0.36, 5.61, 3.10, 0.25, 10, color, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "来源：审计报告 P0 / P1 / P2 结论摘要")

    # 04 / Principles
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "03 / 设计原则", "用五条可验收原则约束全部页面，而非只做换肤", 4)
    principles = [
        ("01", "状态先于内容", "loading / ready / empty / error / unauthorized 明确互斥", RED),
        ("02", "结果先于命令", "目标、执行、同步结果三层表达；失败保留重试路径", ORANGE),
        ("03", "任务先于组件", "总览只回答异常与下一步，高风险设置独立受限", BLUE),
        ("04", "口径先于图表", "范围、样本、更新时间、可比性随图表呈现", PURPLE),
        ("05", "移动先于压缩", "保留异常、执行、确认；记录改为摘要卡与详情", GREEN),
    ]
    for index, (number, title, body, color) in enumerate(principles):
        y = 1.60 + index * 1.04
        add_text(slide, number, 0.83, y + 0.02, 0.62, 0.37, 19, color, True, FONT_NUM)
        add_box(slide, 1.62, y, 10.73, 0.76, PANEL, color, True)
        add_text(slide, title, 1.93, y + 0.16, 2.05, 0.30, 14, color, True)
        add_text(slide, body, 4.12, y + 0.16, 7.82, 0.34, 12, TEXT)
    add_footer(slide, "原则由审计建议归纳，用于后续设计与验收")

    # 05 / IA
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "04 / 信息架构", "五个独立工作区，共享一套连接状态与筛选契约", 5)
    add_box(slide, 0.83, 1.56, 11.70, 0.72, PANEL_2, CYAN, True)
    add_text(slide, "全局壳层：主导航  ·  数据源状态  ·  时间范围  ·  全局搜索  ·  用户权限", 1.10, 1.78, 11.16, 0.28, 13, CYAN, True, align=PP_ALIGN.CENTER)
    workspaces = [
        ("总览", "判断是否需要人工介入", "异常队列\n关键 KPI\n新鲜度", GREEN),
        ("自动化", "确认命令是否真正生效", "目标状态\n实例回执\n计划编排", ORANGE),
        ("分析", "判断趋势是否可比较", "口径说明\n趋势漏斗\n区域分布", PURPLE),
        ("记录", "查找、核验、批量处理", "搜索筛选\n结果列表\n详情审计", BLUE),
        ("设置", "安全编辑高风险配置", "权限边界\n保存差异\n回滚反馈", RED),
    ]
    for index, (name, purpose, content) in enumerate(workspaces):
        x = 0.68 + index * 2.52
        content_color = workspaces[index][3]
        add_box(slide, x, 2.78, 2.22, 2.65, PANEL, content_color, True)
        add_text(slide, name, x + 0.18, 3.03, 1.86, 0.35, 17, content_color, True, align=PP_ALIGN.CENTER)
        add_text(slide, purpose, x + 0.18, 3.58, 1.86, 0.56, 11, TEXT, True, align=PP_ALIGN.CENTER)
        add_text(slide, content, x + 0.24, 4.36, 1.74, 0.72, 10.5, MUTED, align=PP_ALIGN.CENTER)
    add_box(slide, 1.75, 5.82, 9.83, 0.64, "102A2A", GREEN, True)
    add_text(slide, "路由可直达、历史可返回、局部刷新不丢筛选、工作区错误不污染其他页面", 1.98, 6.00, 9.38, 0.28, 12, GREEN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "来源：审计报告页面级优化矩阵")

    # 06 / Overview
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "05 / 总览工作区", "10 秒回答：哪里异常、影响多大、下一步是什么", 6)
    add_screenshot_frame(slide, SCREENSHOTS["overview"], 0.72, 1.52, 8.60, 5.30, "DESKTOP / 总览")
    add_callout(slide, 1, "连接状态显式化", "展示已同步 / 失败 / 未授权、上次成功时间与重试入口。", 9.68, 1.72, 2.90, RED)
    add_callout(slide, 2, "异常优先", "首屏聚合可行动异常，KPI 为判断服务，不做信息堆叠。", 9.68, 3.18, 2.90, ORANGE)
    add_callout(slide, 3, "指标可追溯", "每张指标卡保留范围、更新时间与详情入口。", 9.68, 4.64, 2.90, GREEN)
    add_footer(slide, "原型截图：overview-desktop.png")

    # 07 / Automation
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "06 / 自动化工作区", "把“我点了”与“它生效了”拆成可确认的事件链", 7)
    add_screenshot_frame(slide, SCREENSHOTS["automation"], 0.72, 1.52, 8.60, 5.30, "DESKTOP / 自动化")
    add_callout(slide, 1, "三层状态", "命令目标、当前执行、同步结果各自有文字与颜色。", 9.68, 1.72, 2.90, ORANGE)
    add_callout(slide, 2, "失败可恢复", "待回执、超时、失败均保留查看回执与重试动作。", 9.68, 3.18, 2.90, RED)
    add_callout(slide, 3, "计划可预读", "先给自然语言摘要和下次窗口，再进入高级编排。", 9.68, 4.64, 2.90, CYAN)
    add_footer(slide, "原型截图：automation-desktop.png")

    # 08 / Analytics
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "07 / 分析工作区", "所有趋势图先说明能否比较，再讨论升降", 8)
    add_screenshot_frame(slide, SCREENSHOTS["analytics"], 0.72, 1.52, 8.60, 5.30, "DESKTOP / 分析")
    add_callout(slide, 1, "统一统计范围", "分子与分母绑定同一时间窗口和筛选器。", 9.68, 1.72, 2.90, PURPLE)
    add_callout(slide, 2, "比较契约", "显示样本量、生成时间、denominatorScope 与不可比标记。", 9.68, 3.18, 2.90, BLUE)
    add_callout(slide, 3, "决策导向", "趋势、漏斗、区域围绕同一个管理问题组织。", 9.68, 4.64, 2.90, GREEN)
    add_footer(slide, "原型截图：analytics-desktop.png")

    # 09 / Records
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "08 / 记录工作区", "搜索、结果、批量操作分层，空结果不再吞掉错误", 9)
    add_screenshot_frame(slide, SCREENSHOTS["records"], 0.72, 1.52, 8.60, 5.30, "DESKTOP / 记录")
    add_callout(slide, 1, "状态互斥", "加载失败、未授权、源为空、筛选为空分别呈现。", 9.68, 1.72, 2.90, RED)
    add_callout(slide, 2, "列表稳定", "筛选与选择状态局部保留，刷新不打断当前核验。", 9.68, 3.18, 2.90, BLUE)
    add_callout(slide, 3, "详情可审计", "记录详情承接状态、时间线和可追溯操作。", 9.68, 4.64, 2.90, CYAN)
    add_footer(slide, "原型截图：records-desktop.png")

    # 10 / Settings + mobile
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "09 / 设置与移动", "高风险配置集中管理；移动端保留异常、执行与确认", 10)
    add_screenshot_frame(slide, SCREENSHOTS["settings"], 0.66, 1.52, 7.72, 5.28, "DESKTOP / 设置")
    add_screenshot_frame(slide, SCREENSHOTS["records_mobile"], 8.82, 1.52, 2.55, 5.28, "MOBILE / 记录")
    add_text(slide, "设置", 11.67, 1.76, 0.78, 0.25, 11, RED, True)
    add_text(slide, "独立权限\n保存前差异\n冲突与回滚", 11.67, 2.15, 0.85, 1.10, 10.5, MUTED)
    add_text(slide, "移动", 11.67, 3.58, 0.78, 0.25, 11, GREEN, True)
    add_text(slide, "摘要卡\n详情路径\n可见菜单状态", 11.67, 3.97, 0.85, 1.10, 10.5, MUTED)
    add_text(slide, "验收", 11.67, 5.40, 0.78, 0.25, 11, BLUE, True)
    add_text(slide, "键盘 / 读屏\n320px / 1440px\n200% 缩放", 11.67, 5.79, 0.90, 0.88, 9.5, MUTED)
    add_footer(slide, "原型截图：settings-desktop.png、records-mobile.png")

    # 11 / Priority and conclusion
    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "10 / 实施优先级", "先建立可信底座，再拆任务，最后规模化提升效率", 11)
    roadmap = [
        ("P0", "事实可信", "数据 / 认证 / 控制状态机\n移除“失败即空数组”\n补齐 401 / 网络 / 空数据 E2E", RED),
        ("P1", "运营可控", "拆分五个工作区\n统一命令与回执模型\n并发编辑、弹窗焦点与失败场景", ORANGE),
        ("P2", "效率可达", "字号与对比度提升\n统一 tabs / menu / keyboard 语义\n移动记录视图与漏斗口径", BLUE),
    ]
    for index, (priority, title, body, color) in enumerate(roadmap):
        x = 0.72 + index * 4.18
        add_box(slide, x, 1.65, 3.82, 3.40, PANEL, color, True)
        add_badge(slide, priority, x + 0.28, 1.94, 0.66, color)
        add_text(slide, title, x + 1.12, 1.92, 2.15, 0.34, 17, color, True)
        add_text(slide, body, x + 0.30, 2.72, 3.18, 1.65, 12, TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, 0.72, 5.48, 11.90, 1.00, "102A2A", GREEN, True)
    add_text(slide, "结论", 1.04, 5.78, 0.72, 0.30, 14, GREEN, True)
    add_text(slide, "本次重设计不是一次视觉升级，而是把事实可信、操作可确认和任务可恢复变成产品默认能力。", 1.87, 5.75, 10.28, 0.40, 15, TEXT, True, align=PP_ALIGN.CENTER)
    add_footer(slide, "来源：审计报告“实施优先级”；本页为建议路线，不代表生产已修复")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description="生成仪表盘 UI/UX 重设计管理评审 PPT")
    parser.add_argument("--output", type=Path, default=OUTPUT_PATH)
    args = parser.parse_args()
    validate_inputs()
    output = args.output.resolve()
    build_deck(output)
    print(
        json.dumps(
            {
                "output": str(output),
                "slide_count": 11,
                "screenshots": [str(path.resolve()) for path in SCREENSHOTS.values()],
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
