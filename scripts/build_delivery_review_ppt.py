"""Build a management review deck for the 2026-07-19 delivery funnel.

The script is intentionally read-only against the repository data sources. It
also writes a JSON snapshot so every number in the deck can be independently
checked without opening the PowerPoint file.
"""

from __future__ import annotations

import argparse
import json
import sqlite3
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
REPORT_DATE = "2026-07-19"

BG = "081A26"
PANEL = "102A38"
PANEL_2 = "143444"
TEXT = "F4F8FA"
MUTED = "91AAB5"
CYAN = "40D9C0"
BLUE = "6EB6FF"
ORANGE = "FFB454"
RED = "FF6B6B"
PURPLE = "A995FF"
GREEN = "65D391"
WHITE = "FFFFFF"
FONT = "Microsoft YaHei"
FONT_NUM = "Aptos Display"


def read_jsonl(path: Path) -> list[dict]:
    records: list[dict] = []
    with path.open(encoding="utf-8") as handle:
        for line in handle:
            try:
                records.append(json.loads(line))
            except json.JSONDecodeError:
                continue
    return records


def today_records(records: list[dict]) -> list[dict]:
    return [item for item in records if str(item.get("loggedAt", "")).startswith(REPORT_DATE)]


def parse_dt(value: str) -> datetime | None:
    try:
        return datetime.fromisoformat(value)
    except (TypeError, ValueError):
        return None


def classify_ai_reason(reason: str) -> str:
    text = reason or ""
    if any(key in text for key in ("电工", "电气", "机电", "高低压", "强电", "弱电")):
        return "电气/电工等非目标方向"
    if any(key in text for key in ("英语", "英文", "双语", "日语", "语言", "巴西", "荷兰", "海外")):
        return "语言或海外要求"
    if any(key in text for key in ("经验", "年限", "简历仅", "资历")):
        return "经验/资历不匹配"
    if any(key in text for key in ("学历", "本科", "大专", "专业")):
        return "学历/专业不匹配"
    return "技能或岗位方向不匹配"


def load_snapshot() -> dict:
    actions = today_records(read_jsonl(ROOT / "job_actions.jsonl"))
    all_actions = read_jsonl(ROOT / "job_actions.jsonl")
    ai_today = today_records(read_jsonl(ROOT / "ai_filter_log.jsonl"))
    action_counts = Counter(item.get("action") for item in actions)

    db_path = ROOT / "delivery_state.db"
    connection = sqlite3.connect(f"file:{db_path.as_posix()}?mode=ro", uri=True)
    connection.row_factory = sqlite3.Row
    entered_rows = [
        dict(row)
        for row in connection.execute(
            """
            SELECT account_id, status, claimed_at
            FROM company_deliveries
            WHERE substr(claimed_at, 1, 10) = ?
            ORDER BY claimed_at
            """,
            (REPORT_DATE,),
        )
    ]
    connection.close()

    evaluated = action_counts["job_decision_consumed"]
    entered = len(entered_rows)
    sent = sum(row["status"] == "sent" for row in entered_rows)
    queued = sum(row["status"] == "queued" for row in entered_rows)
    reserved = sum(row["status"] == "reserved" for row in entered_rows)
    failed = sum(row["status"] == "failed_unknown" for row in entered_rows)
    all_time_evaluated = sum(item.get("action") == "job_decision_consumed" for item in all_actions)

    loss_breakdown = [
        {"label": "HR 活跃筛除", "value": action_counts["job_hr_filtered"], "color": RED},
        {
            "label": "领取前未记录/中断",
            "value": evaluated
            - action_counts["job_hr_filtered"]
            - action_counts["job_ai_rejected"]
            - action_counts["job_random_skipped"]
            - action_counts["job_below_threshold"]
            - entered,
            "color": ORANGE,
        },
        {"label": "AI 筛除（HR通过后）", "value": action_counts["job_ai_rejected"], "color": PURPLE},
        {"label": "随机跳过", "value": action_counts["job_random_skipped"], "color": BLUE},
        {"label": "规则阈值", "value": action_counts["job_below_threshold"], "color": MUTED},
    ]

    account_stats: dict[str, dict] = defaultdict(lambda: {"entered": 0, "sent": 0, "limit": 60})
    for row in entered_rows:
        account = row["account_id"] or "未标识账号"
        account_stats[account]["entered"] += 1
        account_stats[account]["sent"] += row["status"] == "sent"
    for account in account_stats:
        account_stats[account]["rate"] = round(account_stats[account]["entered"] / 60 * 100, 1)

    account_events: dict[str, list[datetime]] = defaultdict(list)
    for item in actions:
        timestamp = parse_dt(item.get("loggedAt", ""))
        if timestamp and item.get("accountId"):
            account_events[item["accountId"]].append(timestamp)
    all_times = [parse_dt(item.get("loggedAt", "")) for item in actions]
    all_times = [item for item in all_times if item]
    periods: list[dict] = []
    for account, timestamps in sorted(account_events.items()):
        timestamps.sort()
        if not timestamps:
            continue
        start = previous = timestamps[0]
        for current in timestamps[1:]:
            if (current - previous).total_seconds() > 3 * 3600:
                periods.append({"account": account, "start": start.strftime("%H:%M"), "end": previous.strftime("%H:%M")})
                start = current
            previous = current
        periods.append({"account": account, "start": start.strftime("%H:%M"), "end": previous.strftime("%H:%M")})

    popup_skips = Counter(
        item.get("accountId") or "未标识账号"
        for item in actions
        if item.get("action") == "job_skip" and "浏览器拦截" in (item.get("reason") or "")
    )
    hr_levels = Counter(
        item.get("hrActiveLevel") or "unknown"
        for item in actions
        if item.get("action") == "job_hr_filtered"
    )
    ai_rejections = [item for item in ai_today if item.get("aiPassed") is False]
    ai_categories = Counter(classify_ai_reason(item.get("aiReason", "")) for item in ai_rejections)

    return {
        "date": REPORT_DATE,
        "evaluated": evaluated,
        "entered": entered,
        "queuedEvents": action_counts["greet_queued"],
        "sent": sent,
        "queued": queued,
        "reserved": reserved,
        "failed": failed,
        "allTimeEvaluated": all_time_evaluated,
        "sameDayRate": round(entered / evaluated * 100, 1) if evaluated else 0,
        "sentRate": round(sent / evaluated * 100, 1) if evaluated else 0,
        "displayedRate": round(entered / all_time_evaluated * 100) if all_time_evaluated else 0,
        "lossBreakdown": loss_breakdown,
        "accountStats": dict(account_stats),
        "periods": periods,
        "popupSkips": dict(popup_skips),
        "hrFilteredLevels": dict(hr_levels),
        "aiRejectedTotal": len(ai_rejections),
        "aiCategories": dict(ai_categories),
        "source": [
            "job_actions.jsonl",
            "ai_filter_log.jsonl",
            "delivery_state.db/company_deliveries",
            "dashboard/app.js renderFunnel",
        ],
    }


def save_snapshot(snapshot: dict) -> Path:
    REPORTS.mkdir(parents=True, exist_ok=True)
    path = REPORTS / f"{REPORT_DATE}-delivery-conversion-review-data.json"
    path.write_text(json.dumps(snapshot, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_box(slide, left, top, width, height, fill=PANEL, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.8)
    return shape


def add_text(slide, text, left, top, width, height, size=16, color=TEXT, bold=False, font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Pt(3)
    frame.margin_right = Pt(3)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_title(slide, eyebrow, title, page, snapshot):
    add_text(slide, eyebrow.upper(), 0.72, 0.38, 5.6, 0.25, size=9, color=CYAN, bold=True)
    add_text(slide, title, 0.68, 0.70, 11.5, 0.58, size=26, bold=True)
    add_text(slide, f"数据日期 {snapshot['date']}  ·  管理层复盘", 0.72, 1.30, 6.5, 0.22, size=10, color=MUTED)
    add_text(slide, f"{page:02d} / 07", 12.0, 0.40, 0.65, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, snapshot, note="口径：本地时间；来源见数据快照"):
    add_text(slide, note, 0.72, 7.13, 10.9, 0.18, size=8, color=MUTED)
    add_text(slide, "GOODJOB 运营复盘", 11.0, 7.13, 1.6, 0.18, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_metric(slide, label, value, left, top, width, accent=CYAN, sub=""):
    add_box(slide, left, top, width, 1.05, fill=PANEL, line=PANEL_2, radius=True)
    add_text(slide, label, left + 0.15, top + 0.13, width - 0.3, 0.2, size=10, color=MUTED)
    add_text(slide, value, left + 0.13, top + 0.34, width - 0.26, 0.43, size=25, color=accent, bold=True, font=FONT_NUM)
    if sub:
        add_text(slide, sub, left + 0.15, top + 0.80, width - 0.3, 0.16, size=8, color=MUTED)


def add_bar(slide, label, value, max_value, left, top, width, color, suffix=""):
    add_text(slide, label, left, top, 2.15, 0.23, size=11, color=TEXT)
    track_left = left + 2.15
    track_width = width - 2.95
    add_box(slide, track_left, top + 0.04, track_width, 0.18, fill=PANEL_2, line=PANEL_2, radius=True)
    fill_width = track_width * value / max_value if max_value else 0
    if fill_width:
        add_box(slide, track_left, top + 0.04, fill_width, 0.18, fill=color, line=color, radius=True)
    add_text(slide, f"{value}{suffix}", left + width - 0.72, top - 0.01, 0.72, 0.25, size=11, color=color, bold=True, align=PP_ALIGN.RIGHT, font=FONT_NUM)


def add_arrow(slide, x1, y, x2, color=MUTED):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True


def display_account(account: str) -> str:
    normalized = str(account or "未标识账号")
    if normalized.isdigit() and len(normalized) >= 7:
        return f"账号{normalized[-4:]}"
    return normalized if len(normalized) <= 12 else normalized[-12:]


def build_deck(snapshot: dict, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: executive takeaway.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_text(slide, "投递运营复盘  /  2026.07.19", 0.75, 0.65, 5.5, 0.3, size=11, color=CYAN, bold=True)
    add_text(slide, "不是 3%，\n是 20.1%", 0.70, 1.38, 6.2, 1.45, size=37, color=TEXT, bold=True)
    add_text(slide, "今日进入投递偏少，首要原因是 HR 活跃筛选过严；\n“3%”来自全历史分母，不能代表今日转化。", 0.76, 3.15, 5.9, 0.72, size=16, color=MUTED)
    add_box(slide, 8.05, 1.00, 4.35, 4.75, fill=PANEL, line=PANEL_2, radius=True)
    add_text(slide, "指标纠偏", 8.45, 1.42, 2.0, 0.25, size=11, color=CYAN, bold=True)
    add_text(slide, f"{snapshot['displayedRate']}%", 8.40, 1.85, 3.0, 0.75, size=41, color=RED, bold=True, font=FONT_NUM)
    add_text(slide, "面板显示\n95 ÷ 3,665（全历史）", 8.43, 2.75, 3.15, 0.56, size=13, color=MUTED)
    add_text(slide, f"{snapshot['sameDayRate']}%", 10.65, 3.55, 1.5, 0.55, size=25, color=GREEN, bold=True, font=FONT_NUM, align=PP_ALIGN.RIGHT)
    add_text(slide, "同日真实\n95 ÷ 473", 10.75, 4.18, 1.35, 0.52, size=12, color=TEXT, align=PP_ALIGN.RIGHT)
    add_text(slide, "结论可信度：高（分母问题可由代码直接复核）", 8.43, 5.16, 3.48, 0.3, size=10, color=ORANGE)
    add_metric(slide, "当日评估", str(snapshot["evaluated"]), 0.75, 5.55, 2.20, BLUE)
    add_metric(slide, "进入投递", str(snapshot["entered"]), 3.15, 5.55, 2.20, CYAN)
    add_metric(slide, "已发送", str(snapshot["sent"]), 5.55, 5.55, 2.20, GREEN)
    add_text(slide, "汇报目标：解释量少的真实原因，并把可修正动作排出优先级。", 8.10, 6.55, 4.2, 0.33, size=11, color=TEXT)
    add_footer(slide, snapshot)

    # Slide 2: denominator correction.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_title(slide, "01 / 指标口径", "先纠正指标：分母错了", 2, snapshot)
    add_box(slide, 0.75, 1.82, 5.70, 3.68, fill="2A1D2A", line="5D3349", radius=True)
    add_text(slide, "面板当前算法", 1.08, 2.13, 2.2, 0.25, size=12, color=RED, bold=True)
    add_text(slide, f"{snapshot['entered']} ÷ {snapshot['allTimeEvaluated']}\n= {snapshot['entered'] / snapshot['allTimeEvaluated'] * 100:.1f}%  →  {snapshot['displayedRate']}%", 1.08, 2.75, 4.8, 1.05, size=26, color=TEXT, bold=True, font=FONT_NUM)
    add_text(slide, "renderFunnel 使用 summary.evaluatedJobs\n它是全历史累计值，不随今日筛选切换。", 1.08, 4.30, 4.75, 0.60, size=13, color=MUTED)
    add_box(slide, 6.88, 1.82, 5.70, 3.68, fill="122A2A", line="2D655E", radius=True)
    add_text(slide, "同日复盘算法", 7.21, 2.13, 2.2, 0.25, size=12, color=GREEN, bold=True)
    add_text(slide, f"{snapshot['entered']} ÷ {snapshot['evaluated']}\n= {snapshot['sameDayRate']}%", 7.21, 2.75, 4.8, 1.05, size=29, color=TEXT, bold=True, font=FONT_NUM)
    add_text(slide, "分子、分母都限定在 2026-07-19，\n结果才可用于判断当天筛选效率。", 7.21, 4.30, 4.75, 0.60, size=13, color=MUTED)
    add_text(slide, "判断：3% 是仪表盘展示问题，不是今日系统真实转化。", 0.82, 6.05, 11.8, 0.45, size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, snapshot, "证据：dashboard/app.js renderFunnel；评估量来自 job_actions.jsonl")

    # Slide 3: funnel.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_title(slide, "02 / 当日漏斗", "473 个岗位评估，最终 89 个发送", 3, snapshot)
    stages = [
        ("评估岗位", snapshot["evaluated"], BLUE, 0.85),
        ("进入投递", snapshot["entered"], CYAN, 4.45),
        ("已发送", snapshot["sent"], GREEN, 8.05),
    ]
    for index, (label, value, color, x) in enumerate(stages):
        add_box(slide, x, 2.08, 2.65, 2.15, fill=PANEL, line=color, radius=True)
        add_text(slide, label, x + 0.18, 2.34, 2.25, 0.27, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, str(value), x + 0.20, 2.82, 2.2, 0.70, size=35, color=color, bold=True, font=FONT_NUM, align=PP_ALIGN.CENTER)
        rate = "100%" if index == 0 else f"占评估 {value / snapshot['evaluated'] * 100:.1f}%"
        add_text(slide, rate, x + 0.18, 3.68, 2.25, 0.25, size=11, color=TEXT, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_arrow(slide, x + 2.72, 3.15, x + 3.48, MUTED)
    add_box(slide, 0.85, 4.70, 11.55, 1.00, fill=PANEL, line=PANEL_2, radius=True)
    add_text(slide, "投递状态", 1.12, 4.98, 1.0, 0.2, size=11, color=MUTED)
    status_items = [("已发送", snapshot["sent"], GREEN), ("排队", snapshot["queued"], CYAN), ("预留", snapshot["reserved"], ORANGE), ("异常", snapshot["failed"], RED)]
    for index, (label, value, color) in enumerate(status_items):
        x = 2.35 + index * 2.25
        add_text(slide, f"{label}  {value}", x, 4.91, 1.85, 0.27, size=13, color=color, bold=True)
    add_text(slide, "今日真实进入投递率", 4.30, 6.03, 2.2, 0.25, size=12, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, f"{snapshot['sameDayRate']}%", 6.62, 5.86, 2.4, 0.55, size=29, color=CYAN, bold=True, font=FONT_NUM)
    add_footer(slide, snapshot, "进入投递取 delivery_state.db；已发送/排队/预留/异常为当日状态快照")

    # Slide 4: loss drivers.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_title(slide, "03 / 损耗归因", "最大可确认原因：HR 活跃筛选过严", 4, snapshot)
    add_text(slide, "未进入投递的 378 个岗位", 0.82, 1.78, 3.6, 0.28, size=13, color=MUTED)
    max_loss = max(item["value"] for item in snapshot["lossBreakdown"])
    for index, item in enumerate(snapshot["lossBreakdown"]):
        add_bar(slide, item["label"], item["value"], max_loss, 0.85, 2.26 + index * 0.64, 7.25, item["color"])
        share = item["value"] / (snapshot["evaluated"] - snapshot["entered"]) * 100
        add_text(slide, f"占损耗 {share:.1f}%", 8.55, 2.23 + index * 0.64, 1.3, 0.25, size=10, color=MUTED)
    add_box(slide, 10.25, 1.92, 2.15, 3.85, fill=PANEL, line=PANEL_2, radius=True)
    add_text(slide, "191", 10.48, 2.30, 1.7, 0.64, size=35, color=RED, bold=True, font=FONT_NUM, align=PP_ALIGN.CENTER)
    add_text(slide, "HR 活跃筛除\n占全部损耗 50.5%", 10.45, 3.08, 1.75, 0.63, size=13, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, "当时只允许：\nonline / just_now / today\n\n被筛掉的状态：\nwithin_3_days / this_week / this_month", 10.45, 4.02, 1.75, 1.30, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "132 个“领取前未记录/中断”是日志盲区，不能全部断言为重复岗位。", 0.84, 5.87, 11.3, 0.43, size=14, color=ORANGE, bold=True)
    add_footer(slide, snapshot, "损耗按同日事件与状态快照拆分；132 为剩余项，需补审计事件后再细分")

    # Slide 5: operation.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_title(slide, "04 / 运行效率", "账号没有跑满：第三账号晚启动且弹窗被拦截", 5, snapshot)
    add_text(slide, "按每账号 60 个日上限", 0.82, 1.77, 3.2, 0.26, size=13, color=MUTED)
    account_items = sorted(snapshot["accountStats"].items(), key=lambda pair: pair[1]["entered"], reverse=True)
    for index, (account, values) in enumerate(account_items):
        label = display_account(account)
        add_bar(slide, label, values["entered"], 60, 0.85, 2.25 + index * 0.58, 6.35, [CYAN, BLUE, ORANGE][min(index, 2)])
    add_box(slide, 7.80, 1.82, 4.62, 2.25, fill=PANEL, line=PANEL_2, radius=True)
    add_text(slide, "运行窗口", 8.15, 2.13, 1.4, 0.25, size=12, color=CYAN, bold=True)
    add_text(slide, "02:54–04:57", 8.15, 2.70, 1.8, 0.34, size=18, color=TEXT, bold=True, font=FONT_NUM)
    add_text(slide, "17:40–22:28", 10.25, 2.70, 1.8, 0.34, size=18, color=TEXT, bold=True, font=FONT_NUM)
    add_text(slide, "中间约 12 小时 40 分无评估事件", 8.15, 3.45, 3.65, 0.27, size=11, color=ORANGE)
    add_box(slide, 7.80, 4.38, 4.62, 1.42, fill="2A1D2A", line="5D3349", radius=True)
    add_text(slide, "23 次", 8.15, 4.72, 1.3, 0.45, size=25, color=RED, bold=True, font=FONT_NUM)
    add_text(slide, "第三账号的职位详情被浏览器拦截\n仅进入投递 8 个", 9.45, 4.62, 2.55, 0.65, size=12, color=TEXT)
    add_text(slide, "结论：当日量少既有筛选损耗，也有可用运行时段不足。", 0.84, 5.76, 11.2, 0.36, size=15, color=TEXT, bold=True)
    add_footer(slide, snapshot, "账号进入量取 delivery_state.db；时间段与弹窗拦截取 job_actions.jsonl")

    # Slide 6: evidence levels.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_title(slide, "05 / 证据边界", "哪些原因已经坐实，哪些还需要补日志", 6, snapshot)
    columns = [
        ("已证实", GREEN, ["3% 分母使用全历史评估量", "HR 活跃筛选排除了 191 个", "第三账号晚启动，23 次详情拦截"]),
        ("强推断", ORANGE, ["领取前跳过可能包含重复岗位预检", "长时间停运压缩了可扫描岗位量", "AI 过滤集中在电气/语言/经验不匹配"]),
        ("需补日志", RED, ["132 个评估没有明确后续事件", "无法区分重复、预检不可用与中断", "下一版要记录 skip_reason 与 claim 阶段"]),
    ]
    for index, (heading, color, bullets) in enumerate(columns):
        x = 0.78 + index * 4.18
        add_box(slide, x, 1.95, 3.78, 3.88, fill=PANEL, line=color, radius=True)
        add_text(slide, heading, x + 0.28, 2.25, 2.5, 0.3, size=16, color=color, bold=True)
        for bullet_index, bullet in enumerate(bullets):
            y = 2.92 + bullet_index * 0.78
            add_box(slide, x + 0.30, y + 0.08, 0.12, 0.12, fill=color, line=color, radius=True)
            add_text(slide, bullet, x + 0.54, y, 2.85, 0.46, size=12, color=TEXT)
    add_text(slide, "汇报原则：把“日志未记录”单独列出，避免用未经证实的重复岗位解释全部损耗。", 0.88, 6.20, 11.5, 0.38, size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, snapshot, "证据分级基于源码分支、事件日志和数据库状态；强推断不作为精确归因")

    # Slide 7: actions.
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=BG, line=BG)
    add_title(slide, "06 / 整改优先级", "先修口径，再减少可确认损耗，最后补齐审计", 7, snapshot)
    actions = [
        ("P0  今天就修", RED, "修正漏斗分母\n同日评估量与同日进入量联动\n避免再出现“3%”误判"),
        ("P1  明日验证", ORANGE, "重新评估 HR 活跃范围\n优先验证 within_3_days / this_week\n观察进入率与质量变化"),
        ("P2  本周补齐", CYAN, "为领取前分支统一记录\nskip_reason / precheck / claim\n把 132 个盲区拆成可行动数据"),
    ]
    for index, (heading, color, body) in enumerate(actions):
        x = 0.78 + index * 4.18
        add_box(slide, x, 1.96, 3.78, 3.18, fill=PANEL, line=color, radius=True)
        add_text(slide, heading, x + 0.28, 2.28, 3.15, 0.32, size=16, color=color, bold=True)
        add_text(slide, body, x + 0.28, 3.02, 3.12, 1.35, size=14, color=TEXT)
    add_box(slide, 0.78, 5.58, 11.92, 0.86, fill="122A2A", line="2D655E", radius=True)
    add_text(slide, "明日复盘建议口径：评估岗位 → 同日进入投递 → 已发送；同时单列 HR 筛除、AI 筛除、随机跳过与领取前盲区。", 1.05, 5.82, 11.35, 0.33, size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, snapshot, "建议按 P0/P1/P2 执行；本页为行动建议，不改变历史数据")

    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--data-only", action="store_true")
    parser.add_argument("--output", type=Path, default=REPORTS / f"{REPORT_DATE}-delivery-conversion-review.pptx")
    args = parser.parse_args()
    snapshot = load_snapshot()
    snapshot_path = save_snapshot(snapshot)
    print(json.dumps({"snapshot": str(snapshot_path), **{key: snapshot[key] for key in ("evaluated", "entered", "sent", "sameDayRate", "displayedRate")}}, ensure_ascii=False))
    if not args.data_only:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        build_deck(snapshot, args.output)
        print(json.dumps({"pptx": str(args.output), "slides": 7}, ensure_ascii=False))


if __name__ == "__main__":
    main()
